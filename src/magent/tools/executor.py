"""Built-in tool registry and executor for MagAgent.

Tools are organized into tiers:
  SILENT (0) - read-only, always auto-approved
  AUTO   (1) - shown, auto-approved in balanced/silent mode
  CONFIRM(2) - requires inline y/n prompt in balanced mode
  BLOCK  (3) - always requires explicit typed confirmation
"""

from __future__ import annotations

import asyncio
import base64
import difflib
import fnmatch
import html
import importlib.util
import io
import json
import re
import shlex
import shutil
import sys
import tempfile
import warnings
import zipfile
from contextlib import suppress
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import httpx
from rich.console import Console
from rich.panel import Panel
from rich.prompt import Prompt

from magent.permissions import (
    PermissionResult,
    RiskTier,
    check_permission,
    classify_file_op,
    classify_shell_command,
)
from magent.tool_packs import filter_tool_definitions_for_user
from magent.tools.archive import safe_extract_tar as _safe_extract_tar
from magent.tools.archive import safe_extract_zip as _safe_extract_zip
from magent.tools.registry import tool_def as _def
from magent.tools.registry import validate_tool_args
from magent.tools.types import DEFAULT_TOOL_BUDGETS, READ_FILE_PREVIEW_CHARS, ToolResult

console = Console()

_SEARCH_STOP_WORDS = {
    "about",
    "after",
    "again",
    "also",
    "ancient",
    "based",
    "before",
    "complete",
    "create",
    "current",
    "design",
    "from",
    "history",
    "into",
    "modern",
    "please",
    "research",
    "system",
    "that",
    "the",
    "this",
    "timeline",
    "using",
    "with",
}


def _uses_shell_syntax(command: str) -> bool:
    """Return true when a command needs shell parsing beyond argv splitting."""
    return bool(re.search(r"(\|\||&&|[;\n|<>]|\$\(|\{[^{}\s]+,[^{}]+})", command))


def _first_command_name(command: str) -> str:
    try:
        lexer = shlex.shlex(command, posix=True, punctuation_chars="|&;<>")
        lexer.whitespace_split = True
        for token in lexer:
            if token not in {"|", "||", "&&", ";", ">", ">>", "<"}:
                return Path(token).name.lower()
    except ValueError:
        return ""
    return ""


def _looks_like_read_only_fetch_pipeline(command: str) -> bool:
    """Return true for fetch-and-inspect shell pipelines worth trust-pattern widening."""
    if not re.search(r"\b(curl|wget)\b", command) or "|" not in command:
        return False
    if re.search(r"(?<![<>&])>>?(?![>&])", command) or "<<" in command:
        return False
    try:
        lexer = shlex.shlex(command, posix=True, punctuation_chars="|&;<>")
        lexer.whitespace_split = True
        tokens = list(lexer)
    except ValueError:
        return False
    write_flags = {
        "-d",
        "--data",
        "--data-raw",
        "--data-binary",
        "--data-urlencode",
        "-F",
        "--form",
        "--form-string",
        "-T",
        "--upload-file",
        "-o",
        "--output",
        "-O",
        "--remote-name",
        "--remote-header-name",
        "-X",
        "--request",
    }
    if any(token in write_flags for token in tokens):
        # `-X GET` is technically read-only, but broad approval for arbitrary
        # request overrides is too easy to misapply later.
        return False
    allowed_heads = {"curl", "wget", "grep", "head", "tail", "sed", "awk", "cut", "tr", "od", "wc", "sort", "cat"}
    expect_head = True
    for token in tokens:
        if token in {"|", "||", "&&", ";"}:
            expect_head = True
            continue
        if token in {"2>&1", "1>&2"}:
            continue
        if token in {">", ">>", "<", "2>", "2>>"}:
            return False
        if expect_head:
            if Path(token).name.lower() not in allowed_heads:
                return False
            expect_head = False
    return True


def _effective_shell_timeout(command: str, requested_timeout: int) -> int:
    """Give known package-install commands enough time unless caller overrides."""
    if requested_timeout != 60:
        return requested_timeout
    lowered = f" {command.lower()} "
    package_install_markers = (
        " npm install",
        " npm ci",
        " yarn install",
        " pnpm install",
        " bun install",
    )
    if any(marker in lowered for marker in package_install_markers):
        return 300
    return requested_timeout


def _running_tool_status(tool_name: str, args: dict[str, Any], elapsed: float) -> str:
    label = tool_name
    if tool_name == "run_shell":
        command = str(args.get("command") or "").strip()
        if command:
            label = command[:90]
    return f"Still running {label}... {int(elapsed)}s elapsed"


async def _create_shell_process(command: str, cwd: str | Path):
    """Run shell syntax through bash when available so brace expansion behaves."""
    bash = shutil.which("bash")
    if bash:
        return await asyncio.create_subprocess_exec(
            bash,
            "-lc",
            command,
            cwd=cwd,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
    return await asyncio.create_subprocess_shell(
        command,
        cwd=cwd,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )


async def _create_exec_process(*argv: str, cwd: str | Path):
    return await asyncio.create_subprocess_exec(
        *argv,
        cwd=cwd,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )


_LOW_VALUE_SEARCH_DOMAINS = {
    "facebook.com",
    "instagram.com",
    "pinterest.com",
    "tiktok.com",
    "x.com",
}


def _search_terms(query: str) -> list[str]:
    words = re.findall(r"[a-z0-9][a-z0-9-]{2,}", query.lower())
    terms = [word.strip("-") for word in words if word not in _SEARCH_STOP_WORDS]
    return list(dict.fromkeys(term for term in terms if term))


def _clean_search_result(item: dict[str, Any]) -> dict[str, str]:
    return {
        "title": str(item.get("title") or ""),
        "snippet": str(item.get("body") or item.get("snippet") or ""),
        "url": str(item.get("href") or item.get("url") or ""),
    }


def _is_low_value_search_url(url: str) -> bool:
    host = urlparse(url).netloc.lower().removeprefix("www.")
    return any(host == domain or host.endswith(f".{domain}") for domain in _LOW_VALUE_SEARCH_DOMAINS)


def _prefer_platform_python_command(command: str) -> str:
    """On macOS, prefer the Python 3 command family over ambiguous python/pip."""
    if sys.platform != "darwin":
        return command
    rewritten = re.sub(r"(^|[;&|\n]\s*)pip(?=\s)", r"\1python3 -m pip", command)
    rewritten = re.sub(r"(^|[;&|\n]\s*)python(?=\s)", r"\1python3", rewritten)
    return rewritten


def _shell_native_file_tool_guidance(command: str) -> str:
    """Return guidance when shell is being used for work native file tools should do."""
    scrubbed = re.sub(r"\b[12]?>&[12]\b", "", command)
    scrubbed = re.sub(r"\b[12]?>\s*/dev/null\b", "", scrubbed)
    if "<<" in scrubbed or re.search(r"(?<![<>&])>>?(?![>&])", scrubbed):
        return "Shell redirection/heredocs are not used for file writes. Use write_file or edit_file instead."

    try:
        argv = shlex.split(command)
    except ValueError:
        argv = []
    head = Path(argv[0]).name.lower() if argv else ""
    if head in {"tee", "touch"}:
        return f"`{head}` writes files. Use write_file or edit_file instead."

    lower = command.lower()
    if head in {"python", "python3"} and (
        "open(" in lower
        or ".write_text(" in lower
        or ".write_bytes(" in lower
        or "path(" in lower
    ):
        return "Python shell snippets that write files are disabled. Use write_file or edit_file instead."
    return ""


def _suspicious_write_file_content(path: str, content: str) -> str:
    """Return a reason when a write_file payload is obviously not file content."""
    stripped = (content or "").strip()
    if not stripped:
        return "content is empty"

    target = Path(path)
    path_tokens = {target.name.lower(), str(target).lower()}
    if stripped.lower() in path_tokens:
        return "content is only the target filename/path"

    suffix = target.suffix.lower()
    if suffix in {".html", ".htm"}:
        lower = stripped.lower()
        if len(stripped) < 80 and not any(token in lower for token in ("<html", "<!doctype", "<body", "<section")):
            return "HTML file content is too short and does not look like markup"
    if suffix in {".md", ".markdown"} and len(stripped) < 20 and stripped.lower() == target.stem.lower():
        return "Markdown file content is only the target title"
    return ""


def _rank_search_results(query: str, raw: list[dict[str, Any]], max_results: int) -> list[dict[str, str]]:
    terms = _search_terms(query)
    cleaned = [
        result
        for result in (_clean_search_result(item) for item in raw)
        if result["url"] and not _is_low_value_search_url(result["url"])
    ]
    if not terms:
        return cleaned[:max_results]

    scored: list[tuple[int, int, dict[str, str]]] = []
    for index, result in enumerate(cleaned):
        text = f"{result['title']} {result['snippet']} {result['url']}".lower()
        matched = {term for term in terms if term in text}
        required = 1 if len(terms) == 1 else min(2, len(terms))
        if len(matched) < required:
            continue
        score = len(matched) * 10
        if any(term in result["title"].lower() for term in matched):
            score += 4
        if any(term in result["url"].lower() for term in matched):
            score += 2
        scored.append((score, -index, result))
    scored.sort(reverse=True)
    return [result for _score, _index, result in scored[:max_results]]


def _normalize_document_sections(value: list[dict[str, Any]] | str) -> list[dict[str, Any]]:
    if isinstance(value, str):
        try:
            parsed = json.loads(value)
        except json.JSONDecodeError:
            parsed = [{"title": "Content", "content": value}]
    else:
        parsed = value
    if isinstance(parsed, dict):
        parsed = [parsed]
    sections: list[dict[str, Any]] = []
    for item in parsed if isinstance(parsed, list) else []:
        if isinstance(item, dict):
            bullets = item.get("bullets") or []
            if isinstance(bullets, str):
                bullets = [line.strip("-• ") for line in bullets.splitlines() if line.strip()]
            sections.append(
                {
                    "title": str(item.get("title") or ""),
                    "content": str(item.get("content") or ""),
                    "bullets": [str(bullet) for bullet in bullets if str(bullet).strip()],
                }
            )
        elif str(item).strip():
            sections.append({"title": "", "content": str(item), "bullets": []})
    return sections or [{"title": "Content", "content": "", "bullets": []}]


def _split_paragraphs(text: str) -> list[str]:
    blocks = [block.strip() for block in re.split(r"\n\s*\n", text or "") if block.strip()]
    if blocks:
        return blocks
    return [line.strip() for line in (text or "").splitlines() if line.strip()]


def _normalize_visual_items(value: list[dict[str, Any]] | str | None) -> list[dict[str, Any]]:
    if value is None:
        return []
    if isinstance(value, str):
        try:
            parsed = json.loads(value)
        except json.JSONDecodeError:
            return [{"type": "text", "text": value, "x": 24, "y": 42}]
    else:
        parsed = value
    if isinstance(parsed, dict):
        parsed = [parsed]
    return [dict(item) for item in parsed if isinstance(item, dict)] if isinstance(parsed, list) else []


def _svg_attr(value: Any, default: str = "") -> str:
    text = str(value if value is not None else default)
    return html.escape(text, quote=True)


def _svg_from_elements(width: int, height: int, title: str, elements: list[dict[str, Any]]) -> str:
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}" role="img">',
    ]
    if title:
        parts.append(f"<title>{html.escape(title)}</title>")
    for item in elements:
        kind = str(item.get("type") or item.get("kind") or "text").lower()
        fill = _svg_attr(item.get("fill"), "none")
        stroke = _svg_attr(item.get("stroke"), "#1f2937")
        stroke_width = _svg_attr(item.get("stroke_width") or item.get("stroke-width"), "2")
        if kind in {"rect", "rectangle", "card"}:
            parts.append(
                '<rect '
                f'x="{_svg_attr(item.get("x"), "0")}" y="{_svg_attr(item.get("y"), "0")}" '
                f'width="{_svg_attr(item.get("width") or item.get("w"), "120")}" '
                f'height="{_svg_attr(item.get("height") or item.get("h"), "80")}" '
                f'rx="{_svg_attr(item.get("rx"), "8")}" fill="{fill}" stroke="{stroke}" '
                f'stroke-width="{stroke_width}" />'
            )
        elif kind in {"circle", "ellipse"}:
            if kind == "ellipse" or item.get("rx") or item.get("ry"):
                parts.append(
                    '<ellipse '
                    f'cx="{_svg_attr(item.get("cx") or item.get("x"), "60")}" '
                    f'cy="{_svg_attr(item.get("cy") or item.get("y"), "60")}" '
                    f'rx="{_svg_attr(item.get("rx"), "40")}" ry="{_svg_attr(item.get("ry"), "28")}" '
                    f'fill="{fill}" stroke="{stroke}" stroke-width="{stroke_width}" />'
                )
            else:
                parts.append(
                    '<circle '
                    f'cx="{_svg_attr(item.get("cx") or item.get("x"), "60")}" '
                    f'cy="{_svg_attr(item.get("cy") or item.get("y"), "60")}" '
                    f'r="{_svg_attr(item.get("r") or item.get("radius"), "32")}" '
                    f'fill="{fill}" stroke="{stroke}" stroke-width="{stroke_width}" />'
                )
        elif kind == "line":
            parts.append(
                '<line '
                f'x1="{_svg_attr(item.get("x1"), "0")}" y1="{_svg_attr(item.get("y1"), "0")}" '
                f'x2="{_svg_attr(item.get("x2"), "100")}" y2="{_svg_attr(item.get("y2"), "100")}" '
                f'stroke="{stroke}" stroke-width="{stroke_width}" />'
            )
        elif kind == "path" and item.get("d"):
            parts.append(
                f'<path d="{_svg_attr(item.get("d"))}" fill="{fill}" stroke="{stroke}" '
                f'stroke-width="{stroke_width}" />'
            )
        else:
            text = html.escape(str(item.get("text") or item.get("label") or ""))
            parts.append(
                '<text '
                f'x="{_svg_attr(item.get("x"), "24")}" y="{_svg_attr(item.get("y"), "32")}" '
                f'fill="{_svg_attr(item.get("color") or item.get("fill"), "#111827")}" '
                f'font-size="{_svg_attr(item.get("font_size") or item.get("font-size"), "18")}" '
                f'font-family="{_svg_attr(item.get("font_family") or item.get("font-family"), "system-ui, sans-serif")}">{text}</text>'
            )
    parts.append("</svg>")
    return "\n".join(parts) + "\n"


def _mermaid_from_graph(title: str, nodes: list[dict[str, Any]], edges: list[dict[str, Any]], direction: str) -> str:
    lines = [f"flowchart {direction or 'TD'}"]
    if title:
        lines.append(f"%% {title}")
    known: set[str] = set()
    for index, node in enumerate(nodes, start=1):
        node_id = re.sub(r"\W+", "_", str(node.get("id") or f"n{index}")).strip("_") or f"n{index}"
        label = str(node.get("label") or node.get("title") or node_id)
        known.add(node_id)
        lines.append(f'    {node_id}["{label.replace(chr(34), chr(39))}"]')
    for index, edge in enumerate(edges, start=1):
        source = re.sub(r"\W+", "_", str(edge.get("from") or edge.get("source") or f"n{index}")).strip("_")
        target = re.sub(r"\W+", "_", str(edge.get("to") or edge.get("target") or f"n{index + 1}")).strip("_")
        label = str(edge.get("label") or "").replace('"', "'")
        connector = f'-- "{label}" -->' if label else "-->"
        if source and target:
            if source not in known:
                lines.append(f'    {source}["{source}"]')
                known.add(source)
            if target not in known:
                lines.append(f'    {target}["{target}"]')
                known.add(target)
            lines.append(f"    {source} {connector} {target}")
    return "\n".join(lines) + "\n"


def _pptx_set_background(slide: Any, color: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _pptx_add_text(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int,
    color: Any,
    bold: bool = False,
    align: Any = None,
) -> Any:
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    lines = (text or "").splitlines() or [""]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Aptos"
    return box


def _normalize_tool_args(tool_name: str, args: dict[str, Any]) -> dict[str, Any]:
    """Normalize common provider/model argument aliases before dispatch."""
    normalized = dict(args)
    path_alias_tools = {
        "read_file",
        "read_file_range",
        "outline_file",
        "write_file",
        "create_docx",
        "create_pptx",
        "create_svg",
        "create_diagram",
        "create_image",
        "generate_image",
        "edit_file",
        "delete_file",
        "list_dir",
        "open_file",
        "read_image",
        "browser_screenshot",
    }
    if tool_name in path_alias_tools and "path" not in normalized:
        for alias in ("file_path", "filepath", "filename", "file", "target_path", "output_file"):
            if alias in normalized:
                normalized["path"] = normalized[alias]
                break
    if tool_name == "write_file" and "content" not in normalized:
        for alias in ("contents", "text", "body", "data"):
            if alias in normalized:
                normalized["content"] = normalized[alias]
                break
    if tool_name == "edit_file":
        if "old_str" not in normalized and "old_string" in normalized:
            normalized["old_str"] = normalized["old_string"]
        if "new_str" not in normalized and "new_string" in normalized:
            normalized["new_str"] = normalized["new_string"]
    return normalized


class ToolExecutor:
    """Executes agent tools with integrated permission checking."""

    def __init__(
        self,
        cwd: str,
        permission_mode: str = "balanced",
        allowed_shell_patterns: list[str] | None = None,
        trusted_shell_patterns: list[str] | None = None,
        show_tool_calls: bool = True,
        username: str = "default",
        tool_budgets: dict[str, int] | None = None,
        session_id: str = "manual",
        interactive_permissions: bool = True,
        config: Any | None = None,
    ):
        self.cwd = cwd
        self.permission_mode = permission_mode
        self.allowed_shell_patterns = allowed_shell_patterns or []
        self.trusted_shell_patterns = trusted_shell_patterns or []
        self.session_shell_patterns: list[str] = []
        self.show_tool_calls = show_tool_calls
        self.username = username
        self.tool_budgets = {**DEFAULT_TOOL_BUDGETS, **(tool_budgets or {})}
        self.session_id = session_id
        self.interactive_permissions = interactive_permissions
        self.config = config
        self._active_tasks: set[asyncio.Task[Any]] = set()
        self._active_processes: set[asyncio.subprocess.Process] = set()

    async def cancel_active(self) -> None:
        """Cancel active tool work and terminate subprocesses owned by this executor."""
        tasks = list(self._active_tasks)
        for task in tasks:
            task.cancel()
        for proc in list(self._active_processes):
            with suppress(ProcessLookupError):
                if proc.returncode is None:
                    proc.kill()
            with suppress(Exception):
                await proc.wait()
        if tasks:
            with suppress(Exception):
                await asyncio.wait_for(asyncio.gather(*tasks, return_exceptions=True), timeout=5)
        self._active_processes.clear()
        self._active_tasks.clear()

    def has_active_work(self) -> bool:
        return bool(self._active_tasks or self._active_processes)

    def _checkpoint(self, abs_path: Path, operation: str) -> str:
        try:
            from magent.workbench import create_checkpoint

            item = create_checkpoint(self.username, self.cwd, abs_path, operation, self.session_id)
            return str(item.get("id", ""))
        except Exception:
            return ""

    def _resolve_path(self, path: str) -> Path:
        root = Path(self.cwd).resolve()
        raw = Path(path).expanduser()
        return (raw if raw.is_absolute() else root / raw).resolve(strict=False)

    def _path_tier(self, op: str, path: str) -> tuple[Path, RiskTier]:
        return self._resolve_path(path), classify_file_op(op, path, self.cwd)

    def _log_tool(self, name: str, desc: str, tier: RiskTier) -> None:
        if self.show_tool_calls and tier > RiskTier.SILENT:
            from magent.permissions import TIER_LABELS

            tier_label = TIER_LABELS.get(tier, str(tier))
            console.print(
                f"  [dim]🔧 {name}[/dim] [dim cyan][{tier_label}][/dim cyan] [dim]{desc[:80]}[/dim]"
            )

    def _check_permission(self, action_description: str, tier: RiskTier) -> PermissionResult:
        return check_permission(
            action_description,
            tier,
            self.permission_mode,
            interactive=self.interactive_permissions,
        )

    def _trusted_shell_match(self, command: str) -> bool:
        patterns = [*self.trusted_shell_patterns, *self.session_shell_patterns]
        return any(fnmatch.fnmatch(command.strip(), pattern) for pattern in patterns)

    def _remember_trusted_shell_pattern(self, command: str) -> None:
        try:
            from magent.config import load_user_profile, save_user_profile

            profile = load_user_profile(self.username)
            permissions = profile.setdefault("permissions", {})
            patterns = list(permissions.get("trusted_shell_patterns") or [])
            if command not in patterns:
                patterns.append(command)
            permissions["trusted_shell_patterns"] = patterns
            save_user_profile(self.username, profile)
            self.trusted_shell_patterns = patterns
        except Exception:
            self.session_shell_patterns.append(command)

    def _shell_trust_pattern(self, command: str, tier: RiskTier) -> str:
        if _looks_like_read_only_fetch_pipeline(command):
            head = _first_command_name(command)
            return f"{head} * | *" if head in {"curl", "wget"} else command
        if tier >= RiskTier.CONFIRM:
            return command
        if not re.search(r"(\|\||&&|[;\n|<>]|\$\()", command):
            return command
        try:
            lexer = shlex.shlex(command, posix=True, punctuation_chars="|&;<>")
            lexer.whitespace_split = True
            tokens = list(lexer)
        except ValueError:
            return command
        pieces: list[str] = []
        expect_head = True
        for token in tokens:
            if token in {"|", "||", "&&", ";"}:
                pieces.append(token)
                expect_head = True
                continue
            if token in {">", ">>", "<", "2>", "2>>", "2>&1", "1>&2"}:
                pieces.append(token)
                expect_head = False
                continue
            if expect_head:
                pieces.append(f"{Path(token).name} *")
                expect_head = False
        return " ".join(pieces) if pieces else command

    def _check_shell_permission(self, command: str, tier: RiskTier) -> PermissionResult:
        if self._trusted_shell_match(command):
            return PermissionResult(True, RiskTier.AUTO, "trusted-shell")
        if not self.interactive_permissions or self.permission_mode == "yolo":
            return self._check_permission(f"Run: `{command}`", tier)
        if tier < RiskTier.CONFIRM:
            return PermissionResult(True, tier, "auto")
        title = "[red]⚠ Permission Required[/red]" if tier == RiskTier.BLOCK else "[yellow]Permission[/yellow]"
        border = "red" if tier == RiskTier.BLOCK else "yellow"
        console.print(
            Panel(
                f"[bold]{'High risk shell action' if tier == RiskTier.BLOCK else 'Shell action requires confirmation'}[/bold]\n\n"
                f"Run: `{command}`\n\n"
                "[dim]Choose whether this approval should last once, for this session, or be saved for future sessions.[/dim]",
                title=title,
                border_style=border,
            )
        )
        choice = Prompt.ask(
            "Approve",
            choices=["once", "session", "always", "no", "o", "s", "a", "n"],
            default="once" if tier == RiskTier.CONFIRM else "no",
        ).lower()
        if choice in {"no", "n"}:
            return PermissionResult(False, tier, "user-denied")
        if choice in {"session", "s"}:
            pattern = self._shell_trust_pattern(command, tier)
            if pattern not in self.session_shell_patterns:
                self.session_shell_patterns.append(pattern)
            console.print(f"[dim]Approved for this session; running `{command}`.[/dim]")
            return PermissionResult(True, tier, "user-session-allow")
        if choice in {"always", "a"}:
            pattern = self._shell_trust_pattern(command, tier)
            self._remember_trusted_shell_pattern(pattern)
            console.print(f"[dim]Saved approval for `{pattern}`; running `{command}`.[/dim]")
            return PermissionResult(True, tier, "user-persistent-allow")
        console.print(f"[dim]Approved once; running `{command}`.[/dim]")
        return PermissionResult(True, tier, "user-confirmed")

    def _permission_denied(self, perm: PermissionResult) -> ToolResult:
        error = "Permission required" if perm.reason == "permission-required" else "Permission denied by user"
        return {
            "ok": False,
            "error": error,
            "permission_required": perm.reason == "permission-required",
            "permission_tier": int(perm.tier),
            "permission_reason": perm.reason,
        }

    # ─────────────────────────────────────────────
    # FILE TOOLS
    # ─────────────────────────────────────────────

    async def read_file(self, path: str) -> ToolResult:
        abs_path, tier = self._path_tier("read", path)
        self._log_tool("read_file", str(abs_path), tier)
        perm = self._check_permission(f"Read {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            content = abs_path.read_text(encoding="utf-8", errors="replace")
            total_lines = len(content.splitlines())
            truncated = len(content) > READ_FILE_PREVIEW_CHARS
            if truncated:
                content = (
                    content[:READ_FILE_PREVIEW_CHARS].rstrip()
                    + "\n\n[...file preview truncated; use outline_file or read_file_range...]"
                )
            return {
                "ok": True,
                "content": content,
                "path": str(abs_path),
                "lines": total_lines,
                "truncated": truncated,
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def read_file_range(
        self,
        path: str,
        start_line: int = 1,
        end_line: int | None = None,
    ) -> ToolResult:
        """Read a 1-based inclusive line range from a file."""
        abs_path, tier = self._path_tier("read", path)
        self._log_tool("read_file_range", f"{abs_path}:{start_line}-{end_line or ''}", tier)
        perm = self._check_permission(f"Read {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            lines = abs_path.read_text(encoding="utf-8", errors="replace").splitlines()
            start = max(start_line, 1)
            end = end_line if end_line is not None else min(start + 199, len(lines))
            end = max(start, min(end, len(lines)))
            selected = lines[start - 1 : end]
            numbered = "\n".join(f"{i}: {line}" for i, line in enumerate(selected, start=start))
            return {
                "ok": True,
                "path": str(abs_path),
                "start_line": start,
                "end_line": end,
                "total_lines": len(lines),
                "content": numbered,
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def outline_file(self, path: str, max_symbols: int = 200) -> ToolResult:
        """Return a compact structural outline for a source file."""
        abs_path, tier = self._path_tier("read", path)
        self._log_tool("outline_file", str(abs_path), tier)
        perm = self._check_permission(f"Outline {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            content = abs_path.read_text(encoding="utf-8", errors="replace")
            lines = content.splitlines()
            suffix = abs_path.suffix.lower()
            symbols: list[dict[str, Any]] = []
            if suffix == ".py":
                import ast

                tree = ast.parse(content)
                for node in ast.walk(tree):
                    if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
                        symbols.append(
                            {
                                "kind": "class" if isinstance(node, ast.ClassDef) else "function",
                                "name": node.name,
                                "line": node.lineno,
                            }
                        )
            else:
                import re

                pattern = re.compile(
                    r"^\s*(class|function|def|async def|export function|const|let|var)\s+([A-Za-z_][\w$]*)"
                )
                for lineno, line in enumerate(lines, start=1):
                    match = pattern.match(line)
                    if match:
                        symbols.append(
                            {"kind": match.group(1).strip(), "name": match.group(2), "line": lineno}
                        )
            symbols.sort(key=lambda item: item["line"])
            return {
                "ok": True,
                "path": str(abs_path),
                "lines": len(lines),
                "symbols": symbols[:max_symbols],
                "truncated": len(symbols) > max_symbols,
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def write_file(self, path: str, content: str) -> ToolResult:
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("write_file", str(abs_path), tier)
        suspicious_reason = _suspicious_write_file_content(path, content)
        if suspicious_reason:
            return {
                "ok": False,
                "error": (
                    f"Refused suspicious write_file payload: {suspicious_reason}. "
                    "Call write_file again with the complete intended file contents."
                ),
                "blocked_by": "write-file-content-guard",
                "path": str(abs_path),
                "bytes": len((content or "").encode()),
            }
        perm = self._check_permission(f"Write {len(content)} chars to {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            if abs_path.exists() and abs_path.is_file():
                existing = abs_path.read_text(encoding="utf-8", errors="replace")
                if existing == content:
                    return {
                        "ok": True,
                        "path": str(abs_path),
                        "bytes": len(content.encode()),
                        "unchanged": True,
                    }
            checkpoint_id = self._checkpoint(abs_path, "write_file")
            abs_path.parent.mkdir(parents=True, exist_ok=True)
            abs_path.write_text(content, encoding="utf-8")
            return {
                "ok": True,
                "path": str(abs_path),
                "bytes": len(content.encode()),
                "checkpoint_id": checkpoint_id,
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def create_docx(
        self,
        path: str,
        title: str,
        sections: list[dict[str, Any]] | str,
        subtitle: str = "",
    ) -> ToolResult:
        """Create a Word document from structured sections."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("create_docx", str(abs_path), tier)
        perm = self._check_permission(f"Create Word document {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.shared import Pt, RGBColor

            normalized_sections = _normalize_document_sections(sections)
            checkpoint_id = self._checkpoint(abs_path, "create_docx")
            abs_path.parent.mkdir(parents=True, exist_ok=True)

            doc = Document()
            normal = doc.styles["Normal"]
            normal.font.name = "Aptos"
            normal.font.size = Pt(11)

            heading = doc.add_heading(str(title or abs_path.stem), 0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if subtitle:
                para = doc.add_paragraph(str(subtitle))
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in para.runs:
                    run.italic = True
                    run.font.color.rgb = RGBColor(0x66, 0x55, 0x44)

            for section in normalized_sections:
                section_title = str(section.get("title") or "").strip()
                if section_title:
                    doc.add_heading(section_title, level=1)
                content = str(section.get("content") or "").strip()
                if content:
                    for paragraph in _split_paragraphs(content):
                        doc.add_paragraph(paragraph)
                for bullet in section.get("bullets") or []:
                    text = str(bullet).strip()
                    if text:
                        doc.add_paragraph(text, style="List Bullet")

            doc.save(abs_path)
            return {
                "ok": True,
                "path": str(abs_path),
                "sections": len(normalized_sections),
                "bytes": abs_path.stat().st_size,
                "checkpoint_id": checkpoint_id,
            }
        except ModuleNotFoundError as e:
            return {
                "ok": False,
                "error": f"Missing document dependency `{e.name}`. Install or upgrade MagAgent to include document support.",
                "path": str(abs_path),
            }
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}

    async def create_pptx(
        self,
        path: str,
        title: str,
        slides: list[dict[str, Any]] | str,
        subtitle: str = "",
    ) -> ToolResult:
        """Create a PowerPoint presentation from structured slides."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("create_pptx", str(abs_path), tier)
        perm = self._check_permission(f"Create PowerPoint presentation {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches

            normalized_slides = _normalize_document_sections(slides)
            checkpoint_id = self._checkpoint(abs_path, "create_pptx")
            abs_path.parent.mkdir(parents=True, exist_ok=True)

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            title_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_set_background(title_slide, RGBColor(0xF7, 0xEF, 0xE1))
            _pptx_add_text(
                title_slide,
                str(title or abs_path.stem),
                0.9,
                1.7,
                11.5,
                1.0,
                size=42,
                bold=True,
                color=RGBColor(0x86, 0x46, 0x16),
                align=PP_ALIGN.CENTER,
            )
            if subtitle:
                _pptx_add_text(
                    title_slide,
                    str(subtitle),
                    1.6,
                    3.0,
                    10.2,
                    0.8,
                    size=22,
                    color=RGBColor(0x5F, 0x45, 0x2C),
                    align=PP_ALIGN.CENTER,
                )

            for slide_data in normalized_slides:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _pptx_set_background(slide, RGBColor(0xFF, 0xFB, 0xF4))
                _pptx_add_text(
                    slide,
                    str(slide_data.get("title") or "Slide"),
                    0.6,
                    0.35,
                    12.0,
                    0.65,
                    size=30,
                    bold=True,
                    color=RGBColor(0x86, 0x46, 0x16),
                )
                bullets = [str(item).strip() for item in slide_data.get("bullets") or [] if str(item).strip()]
                content = str(slide_data.get("content") or "").strip()
                if content and not bullets:
                    bullets = _split_paragraphs(content)
                body = "\n".join(f"- {item}" for item in bullets[:8])
                _pptx_add_text(
                    slide,
                    body,
                    0.9,
                    1.35,
                    11.5,
                    5.5,
                    size=20,
                    color=RGBColor(0x35, 0x2B, 0x20),
                )

            prs.save(abs_path)
            return {
                "ok": True,
                "path": str(abs_path),
                "slides": len(prs.slides),
                "bytes": abs_path.stat().st_size,
                "checkpoint_id": checkpoint_id,
            }
        except ModuleNotFoundError as e:
            return {
                "ok": False,
                "error": f"Missing presentation dependency `{e.name}`. Install or upgrade MagAgent to include presentation support.",
                "path": str(abs_path),
            }
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}

    async def create_svg(
        self,
        path: str,
        elements: list[dict[str, Any]] | str,
        title: str = "",
        width: int = 1200,
        height: int = 800,
    ) -> ToolResult:
        """Create an SVG from simple structured vector elements."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("create_svg", str(abs_path), tier)
        perm = self._check_permission(f"Create SVG {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            normalized = _normalize_visual_items(elements)
            content = _svg_from_elements(int(width), int(height), title, normalized)
            checkpoint_id = self._checkpoint(abs_path, "create_svg")
            abs_path.parent.mkdir(parents=True, exist_ok=True)
            abs_path.write_text(content, encoding="utf-8")
            return {
                "ok": True,
                "path": str(abs_path),
                "elements": len(normalized),
                "bytes": len(content.encode()),
                "checkpoint_id": checkpoint_id,
            }
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}

    async def create_diagram(
        self,
        path: str,
        title: str,
        nodes: list[dict[str, Any]] | str,
        edges: list[dict[str, Any]] | str = "",
        direction: str = "TD",
    ) -> ToolResult:
        """Create a Mermaid diagram file from graph nodes and edges."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("create_diagram", str(abs_path), tier)
        perm = self._check_permission(f"Create diagram {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            normalized_nodes = _normalize_visual_items(nodes)
            normalized_edges = _normalize_visual_items(edges)
            diagram = _mermaid_from_graph(title, normalized_nodes, normalized_edges, direction)
            if abs_path.suffix.lower() in {".md", ".markdown"}:
                content = f"# {title or 'Diagram'}\n\n```mermaid\n{diagram}```\n"
            else:
                content = diagram
            checkpoint_id = self._checkpoint(abs_path, "create_diagram")
            abs_path.parent.mkdir(parents=True, exist_ok=True)
            abs_path.write_text(content, encoding="utf-8")
            return {
                "ok": True,
                "path": str(abs_path),
                "nodes": len(normalized_nodes),
                "edges": len(normalized_edges),
                "bytes": len(content.encode()),
                "checkpoint_id": checkpoint_id,
            }
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}

    async def create_image(
        self,
        path: str,
        elements: list[dict[str, Any]] | str,
        title: str = "",
        width: int = 1200,
        height: int = 800,
        background: str = "#ffffff",
    ) -> ToolResult:
        """Create a simple PNG/JPEG image from structured shapes and text."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("create_image", str(abs_path), tier)
        perm = self._check_permission(f"Create image {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            from PIL import Image, ImageDraw, ImageFont

            normalized = _normalize_visual_items(elements)
            image = Image.new("RGB", (int(width), int(height)), background or "#ffffff")
            draw = ImageDraw.Draw(image)
            if title:
                draw.text((24, 24), str(title), fill="#111827", font=ImageFont.load_default())
            for item in normalized:
                kind = str(item.get("type") or item.get("kind") or "text").lower()
                fill = str(item.get("fill") or item.get("color") or "#f3f4f6")
                outline = str(item.get("stroke") or item.get("outline") or "#1f2937")
                x = int(float(item.get("x", item.get("cx", 40))))
                y = int(float(item.get("y", item.get("cy", 40))))
                if kind in {"rect", "rectangle", "card"}:
                    w = int(float(item.get("width", item.get("w", 160))))
                    h = int(float(item.get("height", item.get("h", 90))))
                    draw.rounded_rectangle([x, y, x + w, y + h], radius=int(item.get("radius", 8)), fill=fill, outline=outline)
                elif kind in {"circle", "ellipse"}:
                    rx = int(float(item.get("rx", item.get("r", item.get("radius", 40)))))
                    ry = int(float(item.get("ry", item.get("r", item.get("radius", 40)))))
                    draw.ellipse([x - rx, y - ry, x + rx, y + ry], fill=fill, outline=outline)
                elif kind == "line":
                    draw.line(
                        [
                            int(float(item.get("x1", x))),
                            int(float(item.get("y1", y))),
                            int(float(item.get("x2", x + 100))),
                            int(float(item.get("y2", y + 100))),
                        ],
                        fill=outline,
                        width=int(item.get("stroke_width", item.get("width", 2))),
                    )
                else:
                    draw.text((x, y), str(item.get("text") or item.get("label") or ""), fill=fill, font=ImageFont.load_default())
            checkpoint_id = self._checkpoint(abs_path, "create_image")
            abs_path.parent.mkdir(parents=True, exist_ok=True)
            image.save(abs_path)
            return {
                "ok": True,
                "path": str(abs_path),
                "elements": len(normalized),
                "bytes": abs_path.stat().st_size,
                "checkpoint_id": checkpoint_id,
            }
        except ModuleNotFoundError as e:
            return {
                "ok": False,
                "error": f"Missing image dependency `{e.name}`. Install or upgrade MagAgent to include image support.",
                "path": str(abs_path),
            }
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}

    async def generate_image(
        self,
        path: str,
        prompt: str,
        aspect_ratio: str = "landscape",
        reference_image: str = "",
    ) -> ToolResult:
        """Generate a PNG using the configured image_maker role."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("generate_image", str(abs_path), tier)
        perm = self._check_permission(f"Generate image {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            from magent.config import load_config
            from magent.image_generation import generate_image_with_role

            config = self.config or load_config(self.username)
            result = await generate_image_with_role(
                config,
                prompt,
                abs_path,
                aspect_ratio=aspect_ratio,
                reference_image=reference_image,
            )
            if result.get("ok"):
                result["checkpoint_id"] = self._checkpoint(abs_path, "generate_image")
            return result
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}

    async def edit_file(self, path: str, old_str: str, new_str: str) -> ToolResult:
        abs_path, tier = self._path_tier("edit", path)
        self._log_tool("edit_file", str(abs_path), tier)
        perm = self._check_permission(f"Edit {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            content = abs_path.read_text(encoding="utf-8")
            if old_str not in content:
                return {"ok": False, "error": f"String not found in {path}"}
            checkpoint_id = self._checkpoint(abs_path, "edit_file")
            abs_path.write_text(content.replace(old_str, new_str, 1), encoding="utf-8")
            return {"ok": True, "path": str(abs_path), "checkpoint_id": checkpoint_id}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def delete_file(self, path: str) -> ToolResult:
        abs_path, tier = self._path_tier("delete", path)
        self._log_tool("delete_file", str(abs_path), tier)
        perm = self._check_permission(f"Delete {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            checkpoint_id = self._checkpoint(abs_path, "delete_file")
            if abs_path.is_dir():
                shutil.rmtree(abs_path)
            else:
                abs_path.unlink()
            return {"ok": True, "path": str(abs_path), "checkpoint_id": checkpoint_id}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def magent_docs_search(self, query: str, limit: int = 5) -> ToolResult:
        """Search MagAgent's built-in documentation."""
        from magent.docs import search_docs

        return {"ok": True, "query": query, "results": search_docs(query, limit=limit)}

    async def list_dir(self, path: str = ".") -> ToolResult:
        abs_path, tier = self._path_tier("read", path)
        self._log_tool("list_dir", str(abs_path), tier)
        perm = self._check_permission(f"List {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            entries = []
            for item in sorted(abs_path.iterdir()):
                entries.append(
                    {
                        "name": item.name,
                        "type": "dir" if item.is_dir() else "file",
                        "size": item.stat().st_size if item.is_file() else None,
                    }
                )
            return {"ok": True, "path": str(abs_path), "entries": entries, "count": len(entries)}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def diff_files(self, path_a: str, path_b: str) -> ToolResult:
        """Unified diff between two files."""
        abs_a, tier_a = self._path_tier("read", path_a)
        abs_b, tier_b = self._path_tier("read", path_b)
        tier = max(tier_a, tier_b)
        self._log_tool("diff_files", f"{abs_a} vs {abs_b}", tier)
        perm = self._check_permission(f"Diff {abs_a} and {abs_b}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            a = abs_a.read_text(encoding="utf-8", errors="replace").splitlines(keepends=True)
            b = abs_b.read_text(encoding="utf-8", errors="replace").splitlines(keepends=True)
            diff = list(difflib.unified_diff(a, b, fromfile=path_a, tofile=path_b))
            return {"ok": True, "diff": "".join(diff), "changed": bool(diff)}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def compress(self, source_path: str, output_path: str, format: str = "zip") -> ToolResult:
        """Compress a file or directory into zip or tar.gz."""
        src, src_tier = self._path_tier("read", source_path)
        out, out_tier = self._path_tier("write", output_path)
        tier = max(RiskTier.AUTO, src_tier, out_tier)
        self._log_tool("compress", f"{source_path} → {output_path}", tier)
        perm = self._check_permission(f"Compress {src} → {out}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            if format == "zip":
                with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zf:
                    if src.is_dir():
                        for f in src.rglob("*"):
                            if f.is_file():
                                zf.write(f, f.relative_to(src.parent))
                    else:
                        zf.write(src, src.name)
            else:
                import tarfile

                mode = "w:gz" if format == "tar.gz" else "w:bz2"
                with tarfile.open(out, mode) as tf:
                    tf.add(src, arcname=src.name)
            return {"ok": True, "output": str(out), "size_bytes": out.stat().st_size}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def extract(self, archive_path: str, output_dir: str = ".") -> ToolResult:
        """Extract a zip or tar archive."""
        src, src_tier = self._path_tier("read", archive_path)
        out, out_tier = self._path_tier("write", output_dir)
        tier = max(RiskTier.AUTO, src_tier, out_tier)
        self._log_tool("extract", archive_path, tier)
        perm = self._check_permission(f"Extract {src} to {out}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            out.mkdir(parents=True, exist_ok=True)
            if archive_path.endswith(".zip"):
                with zipfile.ZipFile(src) as zf:
                    names = zf.namelist()
                    _safe_extract_zip(zf, out)
            else:
                import tarfile

                with tarfile.open(src) as tf:
                    names = tf.getnames()
                    _safe_extract_tar(tf, out)
            return {
                "ok": True,
                "output_dir": str(out),
                "extracted": len(names),
                "files": names[:20],
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    # ─────────────────────────────────────────────
    # SHELL & CODE TOOLS
    # ─────────────────────────────────────────────

    async def run_shell(self, command: str, timeout: int = 60) -> ToolResult:
        original_command = command
        command = _prefer_platform_python_command(command)
        timeout = _effective_shell_timeout(command, timeout)
        if command != original_command and self.show_tool_calls:
            console.print(f"[dim]Using macOS Python command: `{command}`[/dim]")
        native_guidance = _shell_native_file_tool_guidance(command)
        if native_guidance:
            self._log_tool("run_shell", command, RiskTier.BLOCK)
            return {
                "ok": False,
                "error": native_guidance,
                "blocked_by": "native-file-tool-policy",
                "recommended_tool": "write_file",
            }
        tier = RiskTier.AUTO if self._trusted_shell_match(command) else classify_shell_command(command, self.allowed_shell_patterns)
        self._log_tool("run_shell", command, tier)
        perm = self._check_shell_permission(command, tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            try:
                argv = shlex.split(command)
            except ValueError as e:
                return {"ok": False, "error": f"Invalid shell syntax: {e}"}
            if not argv:
                return {"ok": False, "error": "Empty command"}
            uses_shell_control = _uses_shell_syntax(command)
            if uses_shell_control:
                proc = await _create_shell_process(command, self.cwd)
            else:
                proc = await _create_exec_process(*argv, cwd=self.cwd)
            self._active_processes.add(proc)
            try:
                stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
            except TimeoutError:
                with suppress(ProcessLookupError):
                    proc.kill()
                with suppress(Exception):
                    await proc.wait()
                return {"ok": False, "error": f"Command timed out after {timeout}s"}
            except asyncio.CancelledError:
                with suppress(ProcessLookupError):
                    proc.kill()
                with suppress(Exception):
                    await proc.wait()
                raise
            finally:
                self._active_processes.discard(proc)
            if self.show_tool_calls and tier >= RiskTier.CONFIRM:
                output_bytes = len(stdout) + len(stderr)
                status = "completed" if proc.returncode == 0 else f"exited {proc.returncode}"
                console.print(f"[dim]Shell command {status}; captured {output_bytes} bytes.[/dim]")
            return {
                "ok": proc.returncode == 0,
                "returncode": proc.returncode,
                "stdout": stdout.decode("utf-8", errors="replace"),
                "stderr": stderr.decode("utf-8", errors="replace"),
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def run_python(self, code: str, timeout: int = 30) -> ToolResult:
        """Execute Python code in an isolated subprocess and capture output."""
        tier = RiskTier.CONFIRM
        self._log_tool("run_python", f"{len(code)} chars of Python", tier)
        perm = self._check_permission("Execute Python code snippet", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False) as f:
                f.write(code)
                tmp_path = f.name
            proc = await _create_exec_process(sys.executable, tmp_path, cwd=self.cwd)
            self._active_processes.add(proc)
            try:
                stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
            except TimeoutError:
                with suppress(ProcessLookupError):
                    proc.kill()
                with suppress(Exception):
                    await proc.wait()
                return {"ok": False, "error": f"Python execution timed out after {timeout}s"}
            except asyncio.CancelledError:
                with suppress(ProcessLookupError):
                    proc.kill()
                with suppress(Exception):
                    await proc.wait()
                raise
            finally:
                self._active_processes.discard(proc)
                Path(tmp_path).unlink(missing_ok=True)
            return {
                "ok": proc.returncode == 0,
                "returncode": proc.returncode,
                "stdout": stdout.decode("utf-8", errors="replace")[:8000],
                "stderr": stderr.decode("utf-8", errors="replace")[:2000],
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def install_package(self, package: str, version: str = "") -> ToolResult:
        """
        Install a Python package via pip, asking user permission first.
        Use this when a task requires an optional dependency (moviepy, pandas, etc).
        """
        tier = RiskTier.CONFIRM
        pkg_spec = f"{package}=={version}" if version else package
        self._log_tool("install_package", pkg_spec, tier)

        # Check if already installed
        if importlib.util.find_spec(package.replace("-", "_").split("[")[0]) is not None:
            return {"ok": True, "already_installed": True, "package": pkg_spec}

        perm = self._check_permission(f"pip install {pkg_spec}  (required for this task)", tier)
        if not perm.approved:
            denied = self._permission_denied(perm)
            denied["package"] = pkg_spec
            return denied

        result = await self.run_shell(
            f"{sys.executable} -m pip install {pkg_spec} --quiet",
            timeout=120,
        )
        if result["ok"]:
            return {"ok": True, "installed": pkg_spec, "already_installed": False}
        return {"ok": False, "error": result.get("stderr", "pip failed"), "package": pkg_spec}

    # ─────────────────────────────────────────────
    # WEB TOOLS
    # ─────────────────────────────────────────────

    async def search_codebase(self, pattern: str, path: str = ".") -> ToolResult:
        abs_path, tier = self._path_tier("read", path)
        self._log_tool("search_codebase", f"{pattern!r} in {abs_path}", tier)
        perm = self._check_permission(f"Search {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        rg = shutil.which("rg") or shutil.which("grep")
        if not rg:
            return {"ok": False, "error": "No search tool found (rg or grep)"}
        cmd = (
            ["rg", "--line-number", "--no-heading", pattern, str(abs_path)]
            if shutil.which("rg")
            else ["grep", "-rn", pattern, str(abs_path)]
        )
        try:
            proc = await _create_exec_process(*cmd, cwd=self.cwd)
            self._active_processes.add(proc)
            try:
                stdout, _ = await asyncio.wait_for(proc.communicate(), timeout=30)
            finally:
                self._active_processes.discard(proc)
            lines = stdout.decode("utf-8", errors="replace").strip().splitlines()
            return {
                "ok": True,
                "matches": lines[:100],
                "truncated": len(lines) > 100,
                "total": len(lines),
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def web_search(self, query: str, max_results: int = 8) -> ToolResult:
        """Search the web using DuckDuckGo/DDGS (no API key required, real results)."""
        tier = RiskTier.AUTO
        self._log_tool("web_search", query, tier)
        perm = self._check_permission(f"Web search: {query}", tier)
        if not perm.approved:
            return self._permission_denied(perm)

        search_errors: list[str] = []
        max_results = max(1, min(int(max_results), 20))

        # Try ddgs first. Fall back to duckduckgo-search for older installs without
        # leaking the upstream rename warning into user sessions.
        try:
            try:
                from ddgs import DDGS  # type: ignore[import-not-found]

                source = "ddgs"
            except ImportError:
                with warnings.catch_warnings():
                    warnings.filterwarnings("ignore", message=".*duckduckgo_search.*ddgs.*")
                    from duckduckgo_search import DDGS  # type: ignore[import-not-found]

                source = "duckduckgo-search"

            with DDGS() as ddgs:
                raw = list(ddgs.text(query, max_results=max_results))
            results = _rank_search_results(query, raw, max_results)
            return {
                "ok": True,
                "query": query,
                "results": results,
                "source": source,
                "raw_count": len(raw),
                "filtered_count": max(0, len(raw) - len(results)),
                "warning": "" if results else "Search returned no relevant results.",
            }
        except Exception as e:
            search_errors.append(str(e))

        # Fallback: DDG instant answer API
        try:
            async with httpx.AsyncClient(timeout=15) as client:
                resp = await client.get(
                    "https://api.duckduckgo.com/",
                    params={"q": query, "format": "json", "no_html": "1"},
                )
                data = resp.json()
                results = []
                if data.get("AbstractText"):
                    results.append(
                        {
                            "title": data.get("Heading", query),
                            "snippet": data.get("AbstractText"),
                            "url": data.get("AbstractURL"),
                        }
                    )
                for r in data.get("RelatedTopics", [])[:6]:
                    if isinstance(r, dict) and r.get("Text"):
                        results.append(
                            {
                                "title": r.get("Text", "")[:80],
                                "snippet": r.get("Text"),
                                "url": r.get("FirstURL"),
                            }
                        )
                ranked = _rank_search_results(query, results, max_results)
                return {
                    "ok": True,
                    "query": query,
                    "results": ranked,
                    "source": "ddg-instant",
                    "raw_count": len(results),
                    "filtered_count": max(0, len(results) - len(ranked)),
                    "warning": "" if ranked else "Search returned no relevant results.",
                    "fallback_errors": search_errors,
                }
        except Exception as e:
            search_errors.append(str(e))
            return {"ok": False, "error": "; ".join(error for error in search_errors if error)}

    async def web_fetch(self, url: str, extract_article: bool = True) -> ToolResult:
        """
        Fetch and read a URL. Uses trafilatura for clean article extraction
        when extract_article=True (default), otherwise strips HTML tags.
        """
        tier = RiskTier.AUTO
        self._log_tool("web_fetch", url, tier)
        perm = self._check_permission(f"Fetch URL: {url}", tier)
        if not perm.approved:
            return self._permission_denied(perm)

        try:
            async with httpx.AsyncClient(timeout=25, follow_redirects=True) as client:
                resp = await client.get(url, headers={"User-Agent": "MagAgent/0.2"})
                raw_html = resp.text

            # Try trafilatura for clean article text
            if extract_article:
                try:
                    import trafilatura

                    text = (
                        trafilatura.extract(raw_html, include_comments=False, include_tables=True)
                        or ""
                    )
                    if text and len(text) > 200:
                        return {
                            "ok": True,
                            "url": str(resp.url),
                            "status": resp.status_code,
                            "content": text[:10000],
                            "truncated": len(text) > 10000,
                            "extractor": "trafilatura",
                        }
                except ImportError:
                    pass

            # Fallback: regex strip HTML
            import re

            text = re.sub(
                r"<script[^>]*>.*?</script>", "", raw_html, flags=re.DOTALL | re.IGNORECASE
            )
            text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r"<[^>]+>", "", text)
            text = re.sub(r"\n{3,}", "\n\n", text).strip()
            return {
                "ok": True,
                "url": str(resp.url),
                "status": resp.status_code,
                "content": text[:8000],
                "truncated": len(text) > 8000,
                "extractor": "html-strip",
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def deep_research(
        self,
        topic: str,
        questions: list[str] | None = None,
        max_sources: int = 6,
        fetch_sources: bool = True,
    ) -> ToolResult:
        """Run a multi-query web research pass with cited evidence packets."""
        tier = RiskTier.AUTO
        self._log_tool("deep_research", topic, tier)
        perm = self._check_permission(f"Deep web research: {topic}", tier)
        if not perm.approved:
            return self._permission_denied(perm)

        research_questions = [q.strip() for q in (questions or []) if str(q).strip()]
        queries = [topic]
        queries.extend(f"{topic} {question}" for question in research_questions[:4])

        seen_urls: set[str] = set()
        sources: list[dict[str, Any]] = []
        search_packets: list[dict[str, Any]] = []
        for query in queries[:5]:
            search = await self.web_search(query, max_results=max(3, min(max_sources, 8)))
            search_packets.append({"query": query, "ok": search.get("ok"), "source": search.get("source")})
            for result in search.get("results", []) if search.get("ok") else []:
                url = str(result.get("url") or "")
                if not url or url in seen_urls:
                    continue
                seen_urls.add(url)
                sources.append(
                    {
                        "title": result.get("title", ""),
                        "url": url,
                        "snippet": result.get("snippet", ""),
                        "query": query,
                    }
                )
                if len(sources) >= max_sources:
                    break
            if len(sources) >= max_sources:
                break

        evidence: list[dict[str, Any]] = []
        for source in sources:
            packet = dict(source)
            if fetch_sources:
                fetched = await self.web_fetch(source["url"], extract_article=True)
                packet["fetch_ok"] = fetched.get("ok", False)
                packet["status"] = fetched.get("status")
                if fetched.get("ok"):
                    content = str(fetched.get("content", ""))
                    packet["excerpt"] = content[:1600].strip()
                    packet["extractor"] = fetched.get("extractor")
                else:
                    packet["fetch_error"] = fetched.get("error", "")
            evidence.append(packet)

        return {
            "ok": True,
            "topic": topic,
            "questions": research_questions,
            "queries": queries[:5],
            "searches": search_packets,
            "source_count": len(evidence),
            "sources": evidence,
            "summary": _research_summary(topic, evidence),
        }

    async def http_request(
        self,
        method: str,
        url: str,
        headers: dict | None = None,
        body: str | dict | None = None,
        timeout: int = 30,
    ) -> ToolResult:
        """Full HTTP client: GET/POST/PUT/PATCH/DELETE with custom headers and body."""
        tier = RiskTier.AUTO
        self._log_tool("http_request", f"{method.upper()} {url}", tier)
        perm = self._check_permission(f"HTTP {method.upper()} {url}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            kwargs: dict[str, Any] = {
                "headers": headers or {},
                "timeout": timeout,
                "follow_redirects": True,
            }
            if body:
                if isinstance(body, dict):
                    kwargs["json"] = body
                else:
                    kwargs["content"] = body.encode() if isinstance(body, str) else body
            async with httpx.AsyncClient() as client:
                resp = await client.request(method.upper(), url, **kwargs)
            try:
                data = resp.json()
            except Exception:
                data = None
            return {
                "ok": resp.is_success,
                "status": resp.status_code,
                "headers": dict(resp.headers),
                "body_text": resp.text[:5000],
                "body_json": data,
                "truncated": len(resp.text) > 5000,
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def browser_snapshot(self, url: str, wait_ms: int = 500) -> ToolResult:
        """Capture title and body text from a page with Playwright."""
        tier = RiskTier.AUTO
        self._log_tool("browser_snapshot", url, tier)
        perm = self._check_permission(f"Browser snapshot: {url}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        from magent.browser import browser_snapshot

        return await browser_snapshot(url, wait_ms=wait_ms)

    async def browser_screenshot(self, url: str, path: str, wait_ms: int = 500) -> ToolResult:
        """Capture a page screenshot with Playwright."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("browser_screenshot", f"{url} -> {abs_path}", tier)
        perm = self._check_permission(f"Browser screenshot: {url}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        from magent.browser import browser_screenshot

        return await browser_screenshot(url, str(abs_path), wait_ms=wait_ms)

    # ─────────────────────────────────────────────
    # DATA TOOLS
    # ─────────────────────────────────────────────

    async def json_query(self, path_or_json: str, query: str) -> ToolResult:
        """
        Run a JMESPath query over a JSON file (if path exists) or a JSON string.
        Example: json_query("data.json", "items[?status=='active'].name")
        """
        self._log_tool("json_query", f"{query!r}", RiskTier.SILENT)
        try:
            import jmespath

            # Decide if it's a file path or raw JSON
            candidate = Path(self.cwd) / path_or_json
            if candidate.exists():
                abs_path, tier = self._path_tier("read", path_or_json)
                perm = self._check_permission(f"Read JSON {abs_path}", tier)
                if not perm.approved:
                    return self._permission_denied(perm)
                data = json.loads(abs_path.read_text())
            else:
                data = json.loads(path_or_json)
            result = jmespath.search(query, data)
            return {"ok": True, "result": result}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    # ─────────────────────────────────────────────
    # SYSTEM TOOLS
    # ─────────────────────────────────────────────

    async def system_info(self) -> ToolResult:
        """Get CPU, RAM, disk usage, OS info, and Python version."""
        self._log_tool("system_info", "fetching system metrics", RiskTier.SILENT)
        try:
            import psutil

            cpu = psutil.cpu_percent(interval=0.5)
            ram = psutil.virtual_memory()
            disk = psutil.disk_usage(self.cwd)
            return {
                "ok": True,
                "cpu_percent": cpu,
                "ram_total_gb": round(ram.total / 1e9, 2),
                "ram_used_gb": round(ram.used / 1e9, 2),
                "ram_percent": ram.percent,
                "disk_total_gb": round(disk.total / 1e9, 2),
                "disk_used_gb": round(disk.used / 1e9, 2),
                "disk_free_gb": round(disk.free / 1e9, 2),
                "platform": sys.platform,
                "python": sys.version.split()[0],
                "cwd": self.cwd,
            }
        except ImportError:
            return {"ok": False, "error": "psutil not installed. Use install_package('psutil')."}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def notify(self, title: str, message: str, urgency: str = "normal") -> ToolResult:
        """Send a desktop notification to alert the user when a long task completes."""
        self._log_tool("notify", title, RiskTier.SILENT)
        try:
            from plyer import notification

            notification.notify(title=title, message=message, timeout=8)
            return {"ok": True, "title": title, "message": message}
        except ImportError:
            pass
        # Fallback to notify-send
        if shutil.which("notify-send"):
            result = await self.run_shell(
                f"notify-send --urgency={urgency} {json.dumps(title)} {json.dumps(message)}"
            )
            return result
        return {"ok": False, "error": "No notification backend available (plyer or notify-send)"}

    async def clipboard_read(self) -> ToolResult:
        """Read the current system clipboard contents."""
        self._log_tool("clipboard_read", "reading clipboard", RiskTier.SILENT)
        try:
            import pyperclip

            text = pyperclip.paste()
            return {"ok": True, "content": text, "length": len(text)}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def clipboard_write(self, text: str) -> ToolResult:
        """Write text to the system clipboard."""
        tier = RiskTier.AUTO
        self._log_tool("clipboard_write", f"{len(text)} chars", tier)
        perm = self._check_permission(f"Write {len(text)} chars to clipboard", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            import pyperclip

            pyperclip.copy(text)
            return {"ok": True, "length": len(text)}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    async def open_file(self, path: str) -> ToolResult:
        """Open a file in its default application (xdg-open on Linux)."""
        abs_path, file_tier = self._path_tier("read", path)
        tier = max(RiskTier.AUTO, file_tier)
        self._log_tool("open_file", str(abs_path), tier)
        perm = self._check_permission(f"Open file: {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        opener = shutil.which("xdg-open") or shutil.which("open")
        if not opener:
            return {"ok": False, "error": "No file opener found (xdg-open / open)"}
        result = await self.run_shell(f"{opener} {shlex.quote(str(abs_path))}")
        return result

    async def read_image(self, path: str) -> ToolResult:
        """Read image metadata and return base64-encoded content for vision models."""
        abs_path, tier = self._path_tier("read", path)
        self._log_tool("read_image", str(abs_path), tier)
        perm = self._check_permission(f"Read image {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            from PIL import Image

            with Image.open(abs_path) as img:
                meta = {
                    "format": img.format,
                    "mode": img.mode,
                    "size": img.size,
                    "width": img.size[0],
                    "height": img.size[1],
                }
                # Resize for embedding if large
                if img.size[0] > 1024 or img.size[1] > 1024:
                    img.thumbnail((1024, 1024))
                buf = io.BytesIO()
                img.save(buf, format=img.format or "PNG")
                b64 = base64.b64encode(buf.getvalue()).decode()
            return {"ok": True, "metadata": meta, "base64": b64, "path": str(abs_path)}
        except ImportError:
            return {"ok": False, "error": "Pillow not installed. Use install_package('Pillow')."}
        except Exception as e:
            return {"ok": False, "error": str(e)}

    # ─────────────────────────────────────────────
    # DATABASE TOOLS (SQLite)
    # ─────────────────────────────────────────────

    async def db_query(
        self, sql: str, params: list | None = None, db_name: str = "default"
    ) -> ToolResult:
        """SELECT from a named user database. Use db_name to target a specific database."""
        self._log_tool("db_query", f"[{db_name}] {sql[:60]}", RiskTier.SILENT)
        from magent.tools.db import db_query

        return db_query(self.username, sql, params, db_name)

    async def db_execute(
        self, sql: str, params: list | None = None, db_name: str = "default"
    ) -> ToolResult:
        """INSERT/UPDATE/DELETE/CREATE TABLE in a named user database."""
        tier = RiskTier.AUTO
        self._log_tool("db_execute", f"[{db_name}] {sql[:60]}", tier)
        perm = self._check_permission(f"DB write [{db_name}]: {sql[:60]}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        from magent.tools.db import db_execute

        return db_execute(self.username, sql, params, db_name)

    async def db_list_tables(self, db_name: str = "default") -> ToolResult:
        """List tables and row counts in a named user database."""
        self._log_tool("db_list_tables", db_name, RiskTier.SILENT)
        from magent.tools.db import db_list_tables

        return db_list_tables(self.username, db_name)

    async def db_schema(self, table: str, db_name: str = "default") -> ToolResult:
        """Show schema (columns, types) for a table in a named database."""
        self._log_tool("db_schema", f"[{db_name}].{table}", RiskTier.SILENT)
        from magent.tools.db import db_schema

        return db_schema(self.username, table, db_name)

    async def db_list_databases(self) -> ToolResult:
        """List all databases created for the current user."""
        self._log_tool("db_list_databases", self.username, RiskTier.SILENT)
        from magent.tools.db import list_databases

        return list_databases(self.username)

    # ─────────────────────────────────────────────
    # GIT TOOLS
    # ─────────────────────────────────────────────

    async def git_op(self, subcommand: str, *args: str) -> ToolResult:
        cmd = f"git {subcommand} {' '.join(args)}"
        return await self.run_shell(cmd)

    # ─────────────────────────────────────────────
    # TOOL DEFINITIONS (OpenAI function-calling format)
    # ─────────────────────────────────────────────

    def get_tool_definitions(self) -> list[dict[str, Any]]:
        definitions = [
            _def(
                "read_file",
                "Read the contents of a file.",
                {"path": ("string", "Relative file path")},
            ),
            _def(
                "read_file_range",
                "Read a 1-based inclusive range of lines from a file.",
                {
                    "path": ("string", "Relative file path"),
                    "start_line": ("integer", "First line to read"),
                    "end_line": ("integer", "Optional final line to read"),
                },
            ),
            _def(
                "outline_file",
                "Return a compact structural outline of a source file.",
                {
                    "path": ("string", "Relative file path"),
                    "max_symbols": ("integer", "Maximum symbols to return"),
                },
            ),
            _def(
                "write_file",
                "Write content to a file (creates or overwrites). Use this for generated pages, docs, scripts, and other file creation.",
                {"path": ("string", "File path"), "content": ("string", "Full file content")},
            ),
            _def(
                "create_docx",
                "Create a Word .docx document from structured sections. Prefer this over generating Python scripts for Word documents.",
                {
                    "path": ("string", "Output .docx file path"),
                    "title": ("string", "Document title"),
                    "sections": ("array", "Sections with title, content, and optional bullets"),
                    "subtitle": ("string", "Optional subtitle"),
                },
            ),
            _def(
                "create_pptx",
                "Create a PowerPoint .pptx presentation from structured slides. Prefer this over generating Python scripts for presentations.",
                {
                    "path": ("string", "Output .pptx file path"),
                    "title": ("string", "Presentation title"),
                    "slides": ("array", "Slides with title, content, and/or bullets"),
                    "subtitle": ("string", "Optional subtitle"),
                },
            ),
            _def(
                "create_svg",
                "Create an SVG vector image from structured shapes, paths, lines, and text. Prefer this over generating SVGs with shell or Python scripts.",
                {
                    "path": ("string", "Output .svg file path"),
                    "elements": ("array", "Visual elements with type, position, fill/stroke, text, etc."),
                    "title": ("string", "Optional accessible SVG title"),
                    "width": ("integer", "Optional width in pixels"),
                    "height": ("integer", "Optional height in pixels"),
                },
            ),
            _def(
                "create_diagram",
                "Create a Mermaid diagram file from structured nodes and edges.",
                {
                    "path": ("string", "Output .mmd or .md file path"),
                    "title": ("string", "Diagram title"),
                    "nodes": ("array", "Nodes with id and label"),
                    "edges": ("array", "Optional edges with from/source, to/target, and optional label"),
                    "direction": ("string", "Optional Mermaid flowchart direction such as TD or LR"),
                },
            ),
            _def(
                "create_image",
                "Create a simple PNG/JPEG image from structured shapes and text using local rendering.",
                {
                    "path": ("string", "Output image path such as .png or .jpg"),
                    "elements": ("array", "Visual elements with type, position, fill/stroke, text, etc."),
                    "title": ("string", "Optional title text"),
                    "width": ("integer", "Optional width in pixels"),
                    "height": ("integer", "Optional height in pixels"),
                    "background": ("string", "Optional background color"),
                },
            ),
            _def(
                "generate_image",
                "Generate an AI-created PNG image through the configured image_maker model role.",
                {
                    "path": ("string", "Output image path such as .png"),
                    "prompt": ("string", "Detailed visual prompt"),
                    "aspect_ratio": ("string", "Optional: landscape, portrait, or square"),
                    "reference_image": ("string", "Optional local path or URL for image editing/reference"),
                },
            ),
            _def(
                "edit_file",
                "Replace an exact string in a file.",
                {
                    "path": ("string", None),
                    "old_str": ("string", "Exact string to replace"),
                    "new_str": ("string", "Replacement string"),
                },
            ),
            _def("delete_file", "Delete a file or directory.", {"path": ("string", None)}),
            _def(
                "list_dir",
                "List contents of a directory.",
                {"path": ("string", "Path (default: .)")},
            ),
            _def(
                "diff_files",
                "Show unified diff between two files.",
                {"path_a": ("string", "First file"), "path_b": ("string", "Second file")},
            ),
            _def(
                "compress",
                "Compress a file or directory to zip or tar.gz.",
                {
                    "source_path": ("string", None),
                    "output_path": ("string", None),
                    "format": ("string", "zip or tar.gz (default: zip)"),
                },
            ),
            _def(
                "extract",
                "Extract a zip or tar archive.",
                {
                    "archive_path": ("string", None),
                    "output_dir": ("string", "Where to extract (default: .)"),
                },
            ),
            _def(
                "run_shell",
                "Run a shell command in the project directory. Do not use for file creation or edits; use write_file/edit_file instead.",
                {"command": ("string", None), "timeout": ("integer", "Seconds (default 60)")},
            ),
            _def(
                "run_python",
                "Execute Python code in an isolated subprocess and return stdout/stderr.",
                {
                    "code": ("string", "Python code to run"),
                    "timeout": ("integer", "Seconds (default 30)"),
                },
            ),
            _def(
                "install_package",
                "Install a Python package via pip (asks user permission first).",
                {
                    "package": ("string", "Package name e.g. moviepy"),
                    "version": ("string", "Optional pinned version"),
                },
            ),
            _def(
                "search_codebase",
                "Search for a pattern in the codebase using ripgrep.",
                {"pattern": ("string", None), "path": ("string", "Directory (default: .)")},
            ),
            _def(
                "web_search",
                "Search the web using DuckDuckGo (real results, no API key).",
                {
                    "query": ("string", None),
                    "max_results": ("integer", "Number of results (default 8)"),
                },
            ),
            _def(
                "web_fetch",
                "Fetch a URL and return clean article text using trafilatura.",
                {
                    "url": ("string", None),
                    "extract_article": ("boolean", "Use trafilatura for clean text (default true)"),
                },
            ),
            _def(
                "deep_research",
                "Run multi-query web research, fetch source pages, and return cited evidence packets.",
                {
                    "topic": ("string", "Research topic or question"),
                    "questions": ("array", "Optional follow-up questions"),
                    "max_sources": ("integer", "Maximum sources to collect (default 6)"),
                    "fetch_sources": ("boolean", "Fetch and excerpt sources (default true)"),
                },
            ),
            _def(
                "http_request",
                "Make any HTTP request (GET/POST/PUT/PATCH/DELETE) with custom headers and body.",
                {
                    "method": ("string", "GET POST PUT PATCH DELETE"),
                    "url": ("string", None),
                    "headers": ("object", "Optional headers dict"),
                    "body": ("string", "Optional body (string or JSON string)"),
                    "timeout": ("integer", "Seconds (default 30)"),
                },
            ),
            _def(
                "browser_snapshot",
                "Capture page title and visible text using Playwright.",
                {
                    "url": ("string", "Page URL"),
                    "wait_ms": ("integer", "Milliseconds to wait after load"),
                },
            ),
            _def(
                "browser_screenshot",
                "Capture a page screenshot using Playwright.",
                {
                    "url": ("string", "Page URL"),
                    "path": ("string", "Output image path"),
                    "wait_ms": ("integer", "Milliseconds to wait after load"),
                },
            ),
            _def(
                "json_query",
                "Run a JMESPath query over a JSON file or JSON string.",
                {
                    "path_or_json": ("string", "File path or raw JSON string"),
                    "query": ("string", "JMESPath expression e.g. items[?active].name"),
                },
            ),
            _def("system_info", "Get CPU, RAM, disk usage, OS info, and Python version.", {}),
            _def(
                "notify",
                "Send a desktop notification to alert the user.",
                {
                    "title": ("string", None),
                    "message": ("string", None),
                    "urgency": ("string", "low/normal/critical"),
                },
            ),
            _def("clipboard_read", "Read the current system clipboard contents.", {}),
            _def(
                "clipboard_write", "Write text to the system clipboard.", {"text": ("string", None)}
            ),
            _def(
                "open_file", "Open a file in its default application.", {"path": ("string", None)}
            ),
            _def(
                "read_image",
                "Read image metadata and base64 encode it for vision models.",
                {"path": ("string", None)},
            ),
            _def(
                "db_query",
                "SELECT from a named SQLite database.",
                {
                    "sql": ("string", "SELECT statement"),
                    "params": ("array", "Optional parameter list"),
                    "db_name": ("string", "Database name (default, project name, or custom)"),
                },
            ),
            _def(
                "db_execute",
                "INSERT/UPDATE/DELETE/CREATE TABLE in a named SQLite database.",
                {
                    "sql": ("string", "SQL statement"),
                    "params": ("array", "Optional parameter list"),
                    "db_name": ("string", "Database name"),
                },
            ),
            _def(
                "db_list_tables",
                "List tables in a named SQLite database.",
                {"db_name": ("string", "Database name (default: 'default')")},
            ),
            _def(
                "db_schema",
                "Show columns and types for a table.",
                {"table": ("string", None), "db_name": ("string", "Database name")},
            ),
            _def(
                "db_list_databases", "List all SQLite databases created for the current user.", {}
            ),
            _def(
                "git_op",
                "Run a git subcommand.",
                {
                    "subcommand": ("string", "e.g. status, diff, add -A, commit -m 'msg'"),
                    "args": ("array", "Optional extra args"),
                },
            ),
            _def(
                "magent_docs_search",
                "Search MagAgent's built-in documentation for command, configuration, and troubleshooting help.",
                {
                    "query": ("string", None),
                    "limit": ("integer", "Number of results (default 5)"),
                },
            ),
        ]
        return filter_tool_definitions_for_user(definitions, self.username)

    def get_tool_definitions_for_message(self, message: str) -> list[dict[str, Any]]:
        """Return a compact relevant tool subset for a user turn."""
        all_defs = self.get_tool_definitions()
        by_name = {item["function"]["name"]: item for item in all_defs}
        text = message.lower()
        selected = {
            "read_file",
            "read_file_range",
            "outline_file",
            "list_dir",
            "search_codebase",
            "run_shell",
            "edit_file",
            "write_file",
            "git_op",
            "system_info",
            "magent_docs_search",
        }
        if any(word in text for word in ("delete", "remove", "clean up", "rename")):
            selected.add("delete_file")
        if any(word in text for word in ("web", "url", "http", "api", "docs", "latest", "search online")):
            selected.update({"web_search", "web_fetch", "http_request"})
        if any(word in text for word in ("research", "compare", "survey", "investigate", "market")):
            selected.update({"web_search", "web_fetch", "deep_research", "http_request"})
        if any(word in text for word in ("browser", "screenshot", "page", "playwright")):
            selected.update({"browser_snapshot", "browser_screenshot"})
        if any(word in text for word in ("json", "csv", "sqlite", "database", "sql", "dataframe", "query")):
            selected.update(
                {
                    "json_query",
                    "db_query",
                    "db_execute",
                    "db_list_tables",
                    "db_schema",
                    "db_list_databases",
                }
            )
        if any(word in text for word in ("image", "screenshot", "photo", "diagram", "vision")):
            selected.update({"read_image", "generate_image"})
        if any(
            word in text
            for word in (
                "diagram",
                "flowchart",
                "mermaid",
                "svg",
                "vector",
                "visual",
                "image",
                "png",
                "jpg",
                "jpeg",
                "illustration",
            )
        ):
            selected.update({"create_diagram", "create_svg", "create_image", "generate_image"})
        if any(word in text for word in ("zip", "archive", "compress", "extract", "tar")):
            selected.update({"compress", "extract"})
        if any(word in text for word in ("clipboard", "notify", "open file", "desktop")):
            selected.update({"clipboard_read", "clipboard_write", "notify", "open_file"})
        if any(word in text for word in ("diff", "compare")):
            selected.add("diff_files")
        if any(word in text for word in ("install", "package", "dependency")):
            selected.add("install_package")
        if any(
            word in text
            for word in (
                "docx",
                "word doc",
                "word document",
                "document",
                "powerpoint",
                "presentation",
                "pptx",
                "slides",
                "slide deck",
            )
        ):
            selected.update({"create_docx", "create_pptx"})
        if len(text.split()) > 120 or any(word in text for word in ("everything", "full access", "all tools")):
            selected.update(by_name)
        return [by_name[name] for name in by_name if name in selected]

    async def dispatch(self, tool_name: str, tool_args: dict[str, Any]) -> ToolResult:
        """Dispatch a tool call by name."""
        a = _normalize_tool_args(tool_name, tool_args)
        raw = bool(a.pop("raw", False))
        definition = next(
            (item for item in self.get_tool_definitions() if item.get("function", {}).get("name") == tool_name),
            None,
        )
        if definition:
            validation = validate_tool_args(definition, a)
            if not validation["ok"]:
                return {
                    "ok": False,
                    "error": f"Missing required arguments for {tool_name}: {', '.join(validation['missing'])}",
                    "missing": validation["missing"],
                }
        dispatch_map: dict[str, Any] = {
            "read_file": lambda: self.read_file(a["path"]),
            "read_file_range": lambda: self.read_file_range(
                a["path"], a.get("start_line", 1), a.get("end_line")
            ),
            "outline_file": lambda: self.outline_file(a["path"], a.get("max_symbols", 200)),
            "write_file": lambda: self.write_file(a["path"], a["content"]),
            "create_docx": lambda: self.create_docx(
                a["path"], a["title"], a["sections"], a.get("subtitle", "")
            ),
            "create_pptx": lambda: self.create_pptx(
                a["path"], a["title"], a["slides"], a.get("subtitle", "")
            ),
            "create_svg": lambda: self.create_svg(
                a["path"],
                a["elements"],
                a.get("title", ""),
                a.get("width", 1200),
                a.get("height", 800),
            ),
            "create_diagram": lambda: self.create_diagram(
                a["path"],
                a["title"],
                a["nodes"],
                a.get("edges", []),
                a.get("direction", "TD"),
            ),
            "create_image": lambda: self.create_image(
                a["path"],
                a["elements"],
                a.get("title", ""),
                a.get("width", 1200),
                a.get("height", 800),
                a.get("background", "#ffffff"),
            ),
            "generate_image": lambda: self.generate_image(
                a["path"],
                a["prompt"],
                a.get("aspect_ratio", "landscape"),
                a.get("reference_image", ""),
            ),
            "edit_file": lambda: self.edit_file(a["path"], a["old_str"], a["new_str"]),
            "delete_file": lambda: self.delete_file(a["path"]),
            "list_dir": lambda: self.list_dir(a.get("path", ".")),
            "diff_files": lambda: self.diff_files(a["path_a"], a["path_b"]),
            "compress": lambda: self.compress(
                a["source_path"], a["output_path"], a.get("format", "zip")
            ),
            "extract": lambda: self.extract(a["archive_path"], a.get("output_dir", ".")),
            "run_shell": lambda: self.run_shell(a["command"], a.get("timeout", 60)),
            "run_python": lambda: self.run_python(a["code"], a.get("timeout", 30)),
            "install_package": lambda: self.install_package(a["package"], a.get("version", "")),
            "search_codebase": lambda: self.search_codebase(a["pattern"], a.get("path", ".")),
            "web_search": lambda: self.web_search(a["query"], a.get("max_results", 8)),
            "web_fetch": lambda: self.web_fetch(a["url"], a.get("extract_article", True)),
            "deep_research": lambda: self.deep_research(
                a["topic"],
                a.get("questions"),
                a.get("max_sources", 6),
                a.get("fetch_sources", True),
            ),
            "http_request": lambda: self.http_request(
                a["method"], a["url"], a.get("headers"), a.get("body"), a.get("timeout", 30)
            ),
            "browser_snapshot": lambda: self.browser_snapshot(a["url"], a.get("wait_ms", 500)),
            "browser_screenshot": lambda: self.browser_screenshot(
                a["url"], a["path"], a.get("wait_ms", 500)
            ),
            "json_query": lambda: self.json_query(a["path_or_json"], a["query"]),
            "system_info": lambda: self.system_info(),
            "notify": lambda: self.notify(a["title"], a["message"], a.get("urgency", "normal")),
            "clipboard_read": lambda: self.clipboard_read(),
            "clipboard_write": lambda: self.clipboard_write(a["text"]),
            "open_file": lambda: self.open_file(a["path"]),
            "read_image": lambda: self.read_image(a["path"]),
            "db_query": lambda: self.db_query(
                a["sql"], a.get("params"), a.get("db_name", "default")
            ),
            "db_execute": lambda: self.db_execute(
                a["sql"], a.get("params"), a.get("db_name", "default")
            ),
            "db_list_tables": lambda: self.db_list_tables(a.get("db_name", "default")),
            "db_schema": lambda: self.db_schema(a["table"], a.get("db_name", "default")),
            "db_list_databases": lambda: self.db_list_databases(),
            "git_op": lambda: self.git_op(a["subcommand"], *a.get("args", [])),
            "magent_docs_search": lambda: self.magent_docs_search(
                a["query"], a.get("limit", 5)
            ),
        }
        fn = dispatch_map.get(tool_name)
        if fn is None:
            return {"ok": False, "error": f"Unknown tool: {tool_name}"}
        try:
            task = asyncio.create_task(fn())
            self._active_tasks.add(task)
            started = asyncio.get_running_loop().time()
            wait_seconds = 6.0
            while True:
                done, _pending = await asyncio.wait({task}, timeout=wait_seconds)
                if done:
                    result = task.result()
                    break
                if self.show_tool_calls:
                    elapsed = asyncio.get_running_loop().time() - started
                    console.print(f"[dim]{_running_tool_status(tool_name, a, elapsed)}[/dim]")
                wait_seconds = 30.0
        except KeyError as e:
            missing = str(e).strip("'")
            return {
                "ok": False,
                "error": f"Missing required argument '{missing}' for tool {tool_name}",
                "tool": tool_name,
                "args": a,
            }
        except asyncio.CancelledError:
            task.cancel()
            await asyncio.gather(task, return_exceptions=True)
            raise
        finally:
            if "task" in locals():
                self._active_tasks.discard(task)
        return result if raw else self._budget_result(tool_name, result)

    def _budget_result(self, tool_name: str, result: ToolResult) -> ToolResult:
        budget = int(self.tool_budgets.get(tool_name, self.tool_budgets.get("default", 8000)))
        if budget <= 0:
            return result
        changed = False
        output = dict(result)
        for key, value in list(output.items()):
            if isinstance(value, str) and len(value) > budget:
                output[key] = value[:budget].rstrip() + (
                    f"\n\n[...{key} truncated at {budget} chars; pass raw=true for full output...]"
                )
                changed = True
            elif isinstance(value, list) and len(json.dumps(value, default=str)) > budget:
                kept = []
                size = 2
                for item in value:
                    item_size = len(json.dumps(item, default=str))
                    if size + item_size > budget:
                        break
                    kept.append(item)
                    size += item_size
                output[key] = kept
                output[f"{key}_truncated"] = True
                changed = True
        if changed:
            output["budgeted"] = True
            output["budget_chars"] = budget
        return output


def _research_summary(topic: str, evidence: list[dict[str, Any]]) -> str:
    """Build a compact non-LLM research summary from source metadata."""
    if not evidence:
        return f"No sources were found for: {topic}"
    lines = [f"Research summary for: {topic}", "", "Sources reviewed:"]
    for index, item in enumerate(evidence, start=1):
        title = item.get("title") or item.get("url") or f"Source {index}"
        url = item.get("url", "")
        snippet = str(item.get("snippet") or item.get("excerpt") or "").replace("\n", " ").strip()
        if len(snippet) > 220:
            snippet = snippet[:220].rstrip() + "..."
        lines.append(f"{index}. {title} - {url}")
        if snippet:
            lines.append(f"   Evidence: {snippet}")
    lines.extend(
        [
            "",
            "Use the source excerpts above for grounded analysis; verify dates and claims before high-stakes decisions.",
        ]
    )
    return "\n".join(lines)
