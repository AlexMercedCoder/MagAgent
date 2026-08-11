"""Document, diagram, and image artifact tools.

The mixin keeps artifact rendering independent from the central tool dispatcher while
preserving ``ToolExecutor`` as MagAgent's stable public facade.
"""

from __future__ import annotations

import html
import json
import re
from pathlib import Path
from typing import Any

from magent.permissions import PermissionResult, RiskTier
from magent.tools.types import ToolResult


def _normalize_document_sections(value: list[dict[str, Any]] | str) -> list[dict[str, Any]]:
    if isinstance(value, str):
        try:
            parsed = json.loads(value)
        except json.JSONDecodeError:
            parsed = [{"title": "Content", "content": value}]
    else:
        parsed = value
    if isinstance(parsed, dict):
        parsed = [parsed]
    sections: list[dict[str, Any]] = []
    for item in parsed if isinstance(parsed, list) else []:
        if isinstance(item, dict):
            bullets = item.get("bullets") or []
            if isinstance(bullets, str):
                bullets = [line.strip("-* \u2022") for line in bullets.splitlines() if line.strip()]
            sections.append(
                {
                    "title": str(item.get("title") or ""),
                    "content": str(item.get("content") or ""),
                    "bullets": [str(bullet) for bullet in bullets if str(bullet).strip()],
                }
            )
        elif str(item).strip():
            sections.append({"title": "", "content": str(item), "bullets": []})
    return sections or [{"title": "Content", "content": "", "bullets": []}]


def _split_paragraphs(text: str) -> list[str]:
    blocks = [block.strip() for block in re.split(r"\n\s*\n", text or "") if block.strip()]
    if blocks:
        return blocks
    return [line.strip() for line in (text or "").splitlines() if line.strip()]


def _normalize_visual_items(value: list[dict[str, Any]] | str | None) -> list[dict[str, Any]]:
    if value is None:
        return []
    if isinstance(value, str):
        try:
            parsed = json.loads(value)
        except json.JSONDecodeError:
            return [{"type": "text", "text": value, "x": 24, "y": 42}]
    else:
        parsed = value
    if isinstance(parsed, dict):
        parsed = [parsed]
    return (
        [dict(item) for item in parsed if isinstance(item, dict)]
        if isinstance(parsed, list)
        else []
    )


def _svg_attr(value: Any, default: str = "") -> str:
    text = str(value if value is not None else default)
    return html.escape(text, quote=True)


def _svg_from_elements(width: int, height: int, title: str, elements: list[dict[str, Any]]) -> str:
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" viewBox="0 0 {width} {height}" role="img">',
    ]
    if title:
        parts.append(f"<title>{html.escape(title)}</title>")
    for item in elements:
        kind = str(item.get("type") or item.get("kind") or "text").lower()
        fill = _svg_attr(item.get("fill"), "none")
        stroke = _svg_attr(item.get("stroke"), "#1f2937")
        stroke_width = _svg_attr(item.get("stroke_width") or item.get("stroke-width"), "2")
        if kind in {"rect", "rectangle", "card"}:
            parts.append(
                "<rect "
                f'x="{_svg_attr(item.get("x"), "0")}" y="{_svg_attr(item.get("y"), "0")}" '
                f'width="{_svg_attr(item.get("width") or item.get("w"), "120")}" '
                f'height="{_svg_attr(item.get("height") or item.get("h"), "80")}" '
                f'rx="{_svg_attr(item.get("rx"), "8")}" fill="{fill}" stroke="{stroke}" '
                f'stroke-width="{stroke_width}" />'
            )
        elif kind in {"circle", "ellipse"}:
            if kind == "ellipse" or item.get("rx") or item.get("ry"):
                parts.append(
                    "<ellipse "
                    f'cx="{_svg_attr(item.get("cx") or item.get("x"), "60")}" '
                    f'cy="{_svg_attr(item.get("cy") or item.get("y"), "60")}" '
                    f'rx="{_svg_attr(item.get("rx"), "40")}" ry="{_svg_attr(item.get("ry"), "28")}" '
                    f'fill="{fill}" stroke="{stroke}" stroke-width="{stroke_width}" />'
                )
            else:
                parts.append(
                    "<circle "
                    f'cx="{_svg_attr(item.get("cx") or item.get("x"), "60")}" '
                    f'cy="{_svg_attr(item.get("cy") or item.get("y"), "60")}" '
                    f'r="{_svg_attr(item.get("r") or item.get("radius"), "32")}" '
                    f'fill="{fill}" stroke="{stroke}" stroke-width="{stroke_width}" />'
                )
        elif kind == "line":
            parts.append(
                "<line "
                f'x1="{_svg_attr(item.get("x1"), "0")}" y1="{_svg_attr(item.get("y1"), "0")}" '
                f'x2="{_svg_attr(item.get("x2"), "100")}" y2="{_svg_attr(item.get("y2"), "100")}" '
                f'stroke="{stroke}" stroke-width="{stroke_width}" />'
            )
        elif kind == "path" and item.get("d"):
            parts.append(
                f'<path d="{_svg_attr(item.get("d"))}" fill="{fill}" stroke="{stroke}" '
                f'stroke-width="{stroke_width}" />'
            )
        else:
            text = html.escape(str(item.get("text") or item.get("label") or ""))
            parts.append(
                "<text "
                f'x="{_svg_attr(item.get("x"), "24")}" y="{_svg_attr(item.get("y"), "32")}" '
                f'fill="{_svg_attr(item.get("color") or item.get("fill"), "#111827")}" '
                f'font-size="{_svg_attr(item.get("font_size") or item.get("font-size"), "18")}" '
                f'font-family="{_svg_attr(item.get("font_family") or item.get("font-family"), "system-ui, sans-serif")}">{text}</text>'
            )
    parts.append("</svg>")
    return "\n".join(parts) + "\n"


def _mermaid_from_graph(
    title: str, nodes: list[dict[str, Any]], edges: list[dict[str, Any]], direction: str
) -> str:
    lines = [f"flowchart {direction or 'TD'}"]
    if title:
        lines.append(f"%% {title}")
    known: set[str] = set()
    for index, node in enumerate(nodes, start=1):
        node_id = re.sub(r"\W+", "_", str(node.get("id") or f"n{index}")).strip("_") or f"n{index}"
        label = str(node.get("label") or node.get("title") or node_id)
        known.add(node_id)
        lines.append(f'    {node_id}["{label.replace(chr(34), chr(39))}"]')
    for index, edge in enumerate(edges, start=1):
        source = re.sub(
            r"\W+", "_", str(edge.get("from") or edge.get("source") or f"n{index}")
        ).strip("_")
        target = re.sub(
            r"\W+", "_", str(edge.get("to") or edge.get("target") or f"n{index + 1}")
        ).strip("_")
        label = str(edge.get("label") or "").replace('"', "'")
        connector = f'-- "{label}" -->' if label else "-->"
        if source and target:
            if source not in known:
                lines.append(f'    {source}["{source}"]')
                known.add(source)
            if target not in known:
                lines.append(f'    {target}["{target}"]')
                known.add(target)
            lines.append(f"    {source} {connector} {target}")
    return "\n".join(lines) + "\n"


def _pptx_set_background(slide: Any, color: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _pptx_rgb(red: int, green: int, blue: int) -> Any:
    from pptx.dml.color import RGBColor

    return RGBColor(red, green, blue)


def _pptx_add_text(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int,
    color: Any,
    bold: bool = False,
    align: Any = None,
) -> Any:
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    lines = (text or "").splitlines() or [""]
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Aptos"
    return box


class ArtifactToolsMixin:
    """Artifact capability implementation mixed into the public tool executor."""

    username: str
    config: Any | None

    def _path_tier(self, op: str, path: str) -> tuple[Path, RiskTier]:
        raise NotImplementedError

    def _log_tool(self, name: str, desc: str, tier: RiskTier) -> None:
        raise NotImplementedError

    def _check_permission(self, action_description: str, tier: RiskTier) -> PermissionResult:
        raise NotImplementedError

    def _permission_denied(self, perm: PermissionResult) -> ToolResult:
        raise NotImplementedError

    def _checkpoint(self, abs_path: Path, operation: str) -> str:
        raise NotImplementedError

    async def create_docx(
        self,
        path: str,
        title: str,
        sections: list[dict[str, Any]] | str,
        subtitle: str = "",
    ) -> ToolResult:
        """Create a Word document from structured sections."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("create_docx", str(abs_path), tier)
        perm = self._check_permission(f"Create Word document {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.shared import Pt, RGBColor

            normalized_sections = _normalize_document_sections(sections)
            checkpoint_id = self._checkpoint(abs_path, "create_docx")
            abs_path.parent.mkdir(parents=True, exist_ok=True)

            doc = Document()
            normal = doc.styles["Normal"]
            normal.font.name = "Aptos"
            normal.font.size = Pt(11)

            heading = doc.add_heading(str(title or abs_path.stem), 0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if subtitle:
                para = doc.add_paragraph(str(subtitle))
                para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for run in para.runs:
                    run.italic = True
                    run.font.color.rgb = RGBColor(0x66, 0x55, 0x44)

            for section in normalized_sections:
                section_title = str(section.get("title") or "").strip()
                if section_title:
                    doc.add_heading(section_title, level=1)
                content = str(section.get("content") or "").strip()
                if content:
                    for paragraph in _split_paragraphs(content):
                        doc.add_paragraph(paragraph)
                for bullet in section.get("bullets") or []:
                    text = str(bullet).strip()
                    if text:
                        doc.add_paragraph(text, style="List Bullet")

            doc.save(str(abs_path))
            return {
                "ok": True,
                "path": str(abs_path),
                "sections": len(normalized_sections),
                "bytes": abs_path.stat().st_size,
                "checkpoint_id": checkpoint_id,
            }
        except ModuleNotFoundError as e:
            return {
                "ok": False,
                "error": f"Missing document dependency `{e.name}`.",
                "install": 'python -m pip install "mag-agent[docs]"',
                "path": str(abs_path),
            }
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}

    async def create_pptx(
        self,
        path: str,
        title: str,
        slides: list[dict[str, Any]] | str,
        subtitle: str = "",
    ) -> ToolResult:
        """Create a PowerPoint presentation from structured slides."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("create_pptx", str(abs_path), tier)
        perm = self._check_permission(f"Create PowerPoint presentation {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            from pptx import Presentation
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches

            normalized_slides = _normalize_document_sections(slides)
            checkpoint_id = self._checkpoint(abs_path, "create_pptx")
            abs_path.parent.mkdir(parents=True, exist_ok=True)

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            title_slide = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_set_background(title_slide, _pptx_rgb(0xF7, 0xEF, 0xE1))
            _pptx_add_text(
                title_slide,
                str(title or abs_path.stem),
                0.9,
                1.7,
                11.5,
                1.0,
                size=42,
                bold=True,
                color=_pptx_rgb(0x86, 0x46, 0x16),
                align=PP_ALIGN.CENTER,
            )
            if subtitle:
                _pptx_add_text(
                    title_slide,
                    str(subtitle),
                    1.6,
                    3.0,
                    10.2,
                    0.8,
                    size=22,
                    color=_pptx_rgb(0x5F, 0x45, 0x2C),
                    align=PP_ALIGN.CENTER,
                )

            for slide_data in normalized_slides:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                _pptx_set_background(slide, _pptx_rgb(0xFF, 0xFB, 0xF4))
                _pptx_add_text(
                    slide,
                    str(slide_data.get("title") or "Slide"),
                    0.6,
                    0.35,
                    12.0,
                    0.65,
                    size=30,
                    bold=True,
                    color=_pptx_rgb(0x86, 0x46, 0x16),
                )
                bullets = [
                    str(item).strip()
                    for item in slide_data.get("bullets") or []
                    if str(item).strip()
                ]
                content = str(slide_data.get("content") or "").strip()
                if content and not bullets:
                    bullets = _split_paragraphs(content)
                body = "\n".join(f"- {item}" for item in bullets[:8])
                _pptx_add_text(
                    slide,
                    body,
                    0.9,
                    1.35,
                    11.5,
                    5.5,
                    size=20,
                    color=_pptx_rgb(0x35, 0x2B, 0x20),
                )

            prs.save(str(abs_path))
            return {
                "ok": True,
                "path": str(abs_path),
                "slides": len(prs.slides),
                "bytes": abs_path.stat().st_size,
                "checkpoint_id": checkpoint_id,
            }
        except ModuleNotFoundError as e:
            return {
                "ok": False,
                "error": f"Missing presentation dependency `{e.name}`.",
                "install": 'python -m pip install "mag-agent[docs]"',
                "path": str(abs_path),
            }
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}

    async def create_svg(
        self,
        path: str,
        elements: list[dict[str, Any]] | str,
        title: str = "",
        width: int = 1200,
        height: int = 800,
    ) -> ToolResult:
        """Create an SVG from simple structured vector elements."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("create_svg", str(abs_path), tier)
        perm = self._check_permission(f"Create SVG {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            normalized = _normalize_visual_items(elements)
            content = _svg_from_elements(int(width), int(height), title, normalized)
            checkpoint_id = self._checkpoint(abs_path, "create_svg")
            abs_path.parent.mkdir(parents=True, exist_ok=True)
            abs_path.write_text(content, encoding="utf-8")
            return {
                "ok": True,
                "path": str(abs_path),
                "elements": len(normalized),
                "bytes": len(content.encode()),
                "checkpoint_id": checkpoint_id,
            }
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}

    async def create_diagram(
        self,
        path: str,
        title: str,
        nodes: list[dict[str, Any]] | str,
        edges: list[dict[str, Any]] | str = "",
        direction: str = "TD",
    ) -> ToolResult:
        """Create a Mermaid diagram file from graph nodes and edges."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("create_diagram", str(abs_path), tier)
        perm = self._check_permission(f"Create diagram {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            normalized_nodes = _normalize_visual_items(nodes)
            normalized_edges = _normalize_visual_items(edges)
            diagram = _mermaid_from_graph(title, normalized_nodes, normalized_edges, direction)
            if abs_path.suffix.lower() in {".md", ".markdown"}:
                content = f"# {title or 'Diagram'}\n\n```mermaid\n{diagram}```\n"
            else:
                content = diagram
            checkpoint_id = self._checkpoint(abs_path, "create_diagram")
            abs_path.parent.mkdir(parents=True, exist_ok=True)
            abs_path.write_text(content, encoding="utf-8")
            return {
                "ok": True,
                "path": str(abs_path),
                "nodes": len(normalized_nodes),
                "edges": len(normalized_edges),
                "bytes": len(content.encode()),
                "checkpoint_id": checkpoint_id,
            }
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}

    async def create_image(
        self,
        path: str,
        elements: list[dict[str, Any]] | str,
        title: str = "",
        width: int = 1200,
        height: int = 800,
        background: str = "#ffffff",
    ) -> ToolResult:
        """Create a simple PNG/JPEG image from structured shapes and text."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("create_image", str(abs_path), tier)
        perm = self._check_permission(f"Create image {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            from PIL import Image, ImageDraw, ImageFont

            normalized = _normalize_visual_items(elements)
            image = Image.new("RGB", (int(width), int(height)), background or "#ffffff")
            draw = ImageDraw.Draw(image)
            if title:
                draw.text((24, 24), str(title), fill="#111827", font=ImageFont.load_default())
            for item in normalized:
                kind = str(item.get("type") or item.get("kind") or "text").lower()
                fill = str(item.get("fill") or item.get("color") or "#f3f4f6")
                outline = str(item.get("stroke") or item.get("outline") or "#1f2937")
                x = int(float(item.get("x", item.get("cx", 40))))
                y = int(float(item.get("y", item.get("cy", 40))))
                if kind in {"rect", "rectangle", "card"}:
                    w = int(float(item.get("width", item.get("w", 160))))
                    h = int(float(item.get("height", item.get("h", 90))))
                    draw.rounded_rectangle(
                        [x, y, x + w, y + h],
                        radius=int(item.get("radius", 8)),
                        fill=fill,
                        outline=outline,
                    )
                elif kind in {"circle", "ellipse"}:
                    rx = int(float(item.get("rx", item.get("r", item.get("radius", 40)))))
                    ry = int(float(item.get("ry", item.get("r", item.get("radius", 40)))))
                    draw.ellipse([x - rx, y - ry, x + rx, y + ry], fill=fill, outline=outline)
                elif kind == "line":
                    draw.line(
                        [
                            int(float(item.get("x1", x))),
                            int(float(item.get("y1", y))),
                            int(float(item.get("x2", x + 100))),
                            int(float(item.get("y2", y + 100))),
                        ],
                        fill=outline,
                        width=int(item.get("stroke_width", item.get("width", 2))),
                    )
                else:
                    draw.text(
                        (x, y),
                        str(item.get("text") or item.get("label") or ""),
                        fill=fill,
                        font=ImageFont.load_default(),
                    )
            checkpoint_id = self._checkpoint(abs_path, "create_image")
            abs_path.parent.mkdir(parents=True, exist_ok=True)
            image.save(abs_path)
            return {
                "ok": True,
                "path": str(abs_path),
                "elements": len(normalized),
                "bytes": abs_path.stat().st_size,
                "checkpoint_id": checkpoint_id,
            }
        except ModuleNotFoundError as e:
            return {
                "ok": False,
                "error": f"Missing image dependency `{e.name}`.",
                "install": 'python -m pip install "mag-agent[media]"',
                "path": str(abs_path),
            }
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}

    async def generate_image(
        self,
        path: str,
        prompt: str,
        aspect_ratio: str = "landscape",
        reference_image: str = "",
    ) -> ToolResult:
        """Generate a PNG using the configured image_maker role."""
        abs_path, tier = self._path_tier("write", path)
        self._log_tool("generate_image", str(abs_path), tier)
        perm = self._check_permission(f"Generate image {abs_path}", tier)
        if not perm.approved:
            return self._permission_denied(perm)
        try:
            from magent.config import load_config
            from magent.image_generation import generate_image_with_role

            config = self.config or load_config(self.username)
            result = await generate_image_with_role(
                config,
                prompt,
                abs_path,
                aspect_ratio=aspect_ratio,
                reference_image=reference_image,
            )
            if result.get("ok"):
                result["checkpoint_id"] = self._checkpoint(abs_path, "generate_image")
            return result
        except Exception as e:
            return {"ok": False, "error": str(e), "path": str(abs_path)}
